"""Script to generate a presentation deck in .pptx format for Darshini.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path: str):
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    NAVY = RGBColor(15, 23, 42)        # #0f172a
    BLUE = RGBColor(30, 58, 138)       # #1e3a8a
    LIGHT_BLUE = RGBColor(59, 130, 246)# #3b82f6
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(100, 116, 139)     # #64748b
    CARD_BG = RGBColor(241, 245, 249)  # #f1f5f9
    BORDER_COLOR = RGBColor(203, 213, 225)

    def add_header(slide, title_text, subtitle_text):
        # Header background
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        banner.fill.solid()
        banner.fill.fore_color.rgb = BLUE
        banner.line.color.rgb = BLUE

        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.7), Inches(0.85))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(191, 219, 254)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    # Title box
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.0))
    tf1 = tbox.text_frame
    p = tf1.paragraphs[0]
    p.text = "🚆 PASHUPATASTRA"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    p = tf1.add_paragraph()
    p.text = "AI-Assisted Railway Maintenance Scheduling & Disruption Recovery"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True

    p = tf1.add_paragraph()
    p.text = "Phase 1 Domain Review & Synthetic Railway Data Architecture"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(148, 163, 184)

    # Info card on title slide
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card1.line.color.rgb = LIGHT_BLUE

    cbox = s1.shapes.add_textbox(Inches(1.2), Inches(4.9), Inches(10.9), Inches(1.6))
    ctf = cbox.text_frame
    p = ctf.paragraphs[0]
    p.text = "WORKSTREAM PRESENTATION FOR TEAM MILESTONE 1 -> PHASE 1"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    p = ctf.add_paragraph()
    p.text = "• Presenter / Domain Lead: Darshini (feature/domain-data)\n• Integration & Backend Lead: Aryan (feature/backend-api)\n• Smart India Hackathon 2026 | Milestone 1 Verified -> Domain Freeze Complete"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 2: Problem & Darshini's Mission
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "1. The Railway Problem & Domain Mission", "Why realistic domain modeling matters for Indian Railways")

    # Card Left: Problem
    c_l = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    c_l.fill.solid()
    c_l.fill.fore_color.rgb = CARD_BG
    c_l.line.color.rgb = BORDER_COLOR

    tb_l = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(5.0))
    tfl = tb_l.text_frame
    tfl.word_wrap = True
    p = tfl.paragraphs[0]
    p.text = "THE PROBLEM WE ARE SOLVING"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE

    bullets = [
        "Tracks are shared by high-speed express trains and freight.",
        "Maintenance requires track 'possessions' (total traffic block).",
        "Available windows are extremely tight (Night & Midday off-peak).",
        "Emergencies happen daily (rail fractures, OHE faults, delays).",
        "Manual section controllers struggle to balance safety, priority, and minimal train disruption."
    ]
    for b in bullets:
        p = tfl.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY

    # Card Right: Darshini's Role
    c_r = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4))
    c_r.fill.solid()
    c_r.fill.fore_color.rgb = CARD_BG
    c_r.line.color.rgb = LIGHT_BLUE

    tb_r = s2.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(5.0))
    tfr = tb_r.text_frame
    tfr.word_wrap = True
    p = tfr.paragraphs[0]
    p.text = "DARSHINI'S MISSION (DOMAIN + DATA)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE

    bullets_r = [
        "Define the mathematical & operational truth of Indian Railways.",
        "Conduct thorough Domain Review (DOMAIN_REVIEW.md).",
        "Protect shared contracts so teammates don't break each other.",
        "Build deterministic synthetic data generator (generator.py).",
        "Provide rich, reproducible corridor fixtures (Scenarios A, B, C) for the entire pipeline."
    ]
    for b in bullets_r:
        p = tfr.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY

    # -------------------------------------------------------------
    # SLIDE 3: System Pipeline & Team Handover
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "2. End-to-End System Pipeline", "How Darshini's data feeds each teammate across the 6 stages")

    steps = [
        ("STAGE 1: DOMAIN & DATA", "Darshini", "Assets, Corridors, Synthetic Scenarios A/B/C", LIGHT_BLUE),
        ("STAGE 2: ML RISK SCORER", "Ayush", "Priority & Risk Scoring based on Asset Defect/Criticality", GRAY),
        ("STAGE 3: CP-SAT OPTIMIZER", "Tyagi", "Optimal Schedule, Committed Block Protection, Infeasibility", GRAY),
        ("STAGE 4: FASTAPI BACKEND", "Aryan", "REST Endpoints: POST /optimize, /health, /disrupt", GRAY),
        ("STAGE 5: OPERATIONS UI", "Archit", "Corridor Timeline Dashboard (Time x Tracks)", GRAY),
        ("STAGE 6: DISRUPTION SIM", "Tirth", "Inject Disruption -> Trigger Recovery Re-solve", GRAY),
    ]

    top = 1.5
    for i, (title, owner, desc, col) in enumerate(steps):
        c = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Inches(0.82))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(239, 246, 255) if i == 0 else CARD_BG
        c.line.color.rgb = col

        tb = s3.shapes.add_textbox(Inches(1.0), Inches(top + 0.08), Inches(11.3), Inches(0.7))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title} — [{owner}]"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BLUE if i == 0 else NAVY

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        top += 0.92

    # -------------------------------------------------------------
    # SLIDE 4: The 5 Core Domain Freeze Decisions
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "3. Domain Review: 5 Core Decisions for Domain Freeze", "Agreed railway architecture preventing downstream integration breakages")

    decisions = [
        ("1. Track Modeling", "Keep string track_id ('UP-1', 'DOWN-1') for zero solver breaking changes. Attach rich corridor metadata in data layer."),
        ("2. Work Types", "Standardized 6 Indian Railways categories: Track Renewal, Tamping, OHE, Signalling, Inspection, Emergency Repair."),
        ("3. Machinery Constraints", "Model heavy track machines via mutual_exclusion_group (e.g. BCM_01, CSM_01). Prevents solver RCPSP bloat & timeouts."),
        ("4. Train Timetable", "Train traffic represented as 24-hour discrete Possession Windows (Night, Midday, Evening). Blocks must fit inside windows."),
        ("5. Disruption Scenarios", "Defined 4 deterministic disruptions (Track Unavailable, Emergency Injection, Possession Curtailment, Machine Breakdown).")
    ]

    top = 1.5
    for title, desc in decisions:
        c = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Inches(0.95))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = BORDER_COLOR

        tb = s4.shapes.add_textbox(Inches(1.0), Inches(top + 0.08), Inches(11.3), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BLUE

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10.5)
        p.font.color.rgb = NAVY
        top += 1.1

    # -------------------------------------------------------------
    # SLIDE 5: Work Categories & Possession Windows Table
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "4. Maintenance Categories & Operational Windows", "Operational characteristics for Tyagi (Optimizer) & Ayush (Scorer)")

    # Left: Work Categories Summary
    c_l = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    c_l.fill.solid()
    c_l.fill.fore_color.rgb = CARD_BG
    c_l.line.color.rgb = BORDER_COLOR

    tb_l = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(5.0))
    tfl = tb_l.text_frame
    p = tfl.paragraphs[0]
    p.text = "6 MAINTENANCE CATEGORIES (WorkType)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE

    cats = [
        "1. TRACK_RENEWAL (180-240m): Heavy BCM machine, CTR/TSR",
        "2. BALLAST_TAMPING (90-150m): Tie-tamper CSM, packing",
        "3. OHE_MAINTENANCE (75-120m): Tower wagon, power block",
        "4. SIGNALLING_INTERLOCK (60-120m): Point machines & circuits",
        "5. ROUTINE_INSPECTION (45-90m): USFD flaw detection",
        "6. EMERGENCY_REPAIR (90-150m): Fracture repair (Priority 0.99)"
    ]
    for c in cats:
        p = tfl.add_paragraph()
        p.text = "• " + c
        p.font.size = Pt(10.5)
        p.font.color.rgb = NAVY

    # Right: Possession Windows
    c_r = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4))
    c_r.fill.solid()
    c_r.fill.fore_color.rgb = CARD_BG
    c_r.line.color.rgb = BORDER_COLOR

    tb_r = s5.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(5.0))
    tfr = tb_r.text_frame
    p = tfr.paragraphs[0]
    p.text = "24-HOUR POSSESSION SLOTS (1440 MIN)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE

    wins = [
        "NIGHT TRAFFIC BLOCK (00:30 - 05:00 | 0030-0300 min)\n  Primary heavy renewals, tamping, and track machines.",
        "MIDDAY MAINTENANCE SLOT (11:30 - 14:30 | 0690-0870 min)\n  Inspections, OHE isolations, point adjustments.",
        "EVENING OFF-PEAK SLOT (21:00 - 23:30 | 1260-1410 min)\n  Signalling, track circuits, minor urgent repairs.",
        "HEADWAY SAFETY BUFFER: 15 minutes between consecutive blocks on the same track."
    ]
    for w in wins:
        p = tfr.add_paragraph()
        p.text = "• " + w
        p.font.size = Pt(10)
        p.font.color.rgb = NAVY

    # -------------------------------------------------------------
    # SLIDE 6: The 3 Synthetic Scenarios
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "5. Synthetic Scenario Datasets (Generator Output)", "Deterministic, seed-driven corridor benchmarks exported as JSON fixtures")

    scenarios = [
        ("SCENARIO A: BASELINE MAINLINE (corridor_a_blocks.json)", "seed=42 | 2 Tracks (UP-1, DOWN-1) | 12 Block Candidates", "• 100% clean feasible schedule.\n• Precedence chain: BLK-003 (Renewal) -> BLK-004 (Tamping).\n• 2 shared machine groups (BCM_MACHINE_01, CSM_TAMPER_01).\n• Perfect starting fixture for Tyagi's CP-SAT solver and Archit's UI."),
        ("SCENARIO B: HIGH-DENSITY QUADRUPLE (corridor_b_dense.json)", "seed=101 | 4 Tracks (UP-1, UP-2, DOWN-1, DOWN-2) | 24 Block Candidates", "• Simulates 4-track trunk corridor (New Delhi - Kanpur).\n• 4 competing machine groups under tight possession windows.\n• Tests CP-SAT solver trade-offs, capacity saturation, and priority ranking."),
        ("SCENARIO C: DISRUPTION & RE-OPTIMIZATION (corridor_c_disrupted.json)", "seed=999 | 2 Tracks (UP-1, DOWN-1) | 14 Block Candidates", "• Pre-configured with 2 COMMITTED blocks on DOWN-1 (01:00-03:15).\n• Ready for Tirth's simulation to inject TRACK_UNAVAILABLE on UP-1 (02:00-05:00).\n• Proves solver re-optimizes remaining blocks while protecting committed work.")
    ]

    top = 1.5
    for title, subtitle, details in scenarios:
        c = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Inches(1.65))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = LIGHT_BLUE

        tb = s6.shapes.add_textbox(Inches(1.0), Inches(top + 0.08), Inches(11.3), Inches(1.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = BLUE

        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = GRAY

        p = tf.add_paragraph()
        p.text = details
        p.font.size = Pt(9.5)
        p.font.color.rgb = NAVY
        top += 1.85

    # -------------------------------------------------------------
    # SLIDE 7: Deliverables, Test Proof & Handover
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "6. Deliverables, Test Results & Teammate Handover", "All Phase 1 goals verified and ready for team freeze")

    # Left: Test Results
    c_l = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    c_l.fill.solid()
    c_l.fill.fore_color.rgb = CARD_BG
    c_l.line.color.rgb = RGBColor(34, 197, 94) # Green

    tb_l = s7.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(5.0))
    tfl = tb_l.text_frame
    p = tfl.paragraphs[0]
    p.text = "10/10 AUTOMATED TESTS PASSED (0.049s)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 101, 52)

    tests = [
        "test_seed_reproducibility [PASSED]",
        "test_track_assignment_validity [PASSED]",
        "test_time_window_consistency [PASSED]",
        "test_dependency_integrity [PASSED]",
        "test_work_types_are_valid_enums [PASSED]",
        "test_score_boundaries [0.0, 1.0] [PASSED]",
        "test_scenario_b_dense_properties [PASSED]",
        "test_scenario_c_committed_blocks [PASSED]",
        "test_json_fixtures_validity [PASSED]",
        "test_different_seeds_produce_variation [PASSED]"
    ]
    for t in tests:
        p = tfl.add_paragraph()
        p.text = "✓ " + t
        p.font.size = Pt(9.5)
        p.font.color.rgb = NAVY

    # Right: Teammate Handover
    c_r = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4))
    c_r.fill.solid()
    c_r.fill.fore_color.rgb = CARD_BG
    c_r.line.color.rgb = LIGHT_BLUE

    tb_r = s7.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(5.0))
    tfr = tb_r.text_frame
    p = tfr.paragraphs[0]
    p.text = "EXACT HANDOVER FOR TEAMMATES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE

    handovers = [
        "Aryan (Backend): DOMAIN_REVIEW.md is complete. Contracts in contracts/schemas.py are frozen & locked.",
        "Tyagi (Optimizer): Run your solver on corridor_a_blocks.json. Machine exclusions & dependencies ready.",
        "Ayush (ML Scorer): Asset models in models.py have criticality, condition, and defect attributes ready.",
        "Tirth (Simulation): Scenario C has committed blocks ready for your TRACK_UNAVAILABLE test.",
        "Archit (Frontend): Load corridor_a_blocks.json immediately to build the timeline UI."
    ]
    for h in handovers:
        p = tfr.add_paragraph()
        p.text = "• " + h
        p.font.size = Pt(9.5)
        p.font.color.rgb = NAVY

    prs.save(output_path)
    print(f"PowerPoint presentation created at: {output_path}")

if __name__ == "__main__":
    out_dir = r"C:\Users\Hp\.gemini\antigravity\scratch\pashupatastra"
    create_deck(os.path.join(out_dir, "Pashupatastra_Darshini_Team_Presentation.pptx"))
